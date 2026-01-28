"""
Document Processing Module
Handles loading and processing of various document types
"""
from pathlib import Path
from typing import List, Optional, Dict, Any
import os

# Document parsers
import PyPDF2
import docx
from pptx import Presentation

from config import settings
from utils.logger import log


class Document:
    """Simple document class to hold content and metadata"""
    
    def __init__(self, page_content: str, metadata: Optional[Dict[str, Any]] = None):
        self.page_content = page_content
        self.metadata = metadata or {}
    
    def __repr__(self):
        return f"Document(page_content='{self.page_content[:50]}...', metadata={self.metadata})"


class DocumentProcessor:
    """Process and chunk documents for RAG pipeline"""
    
    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.txt', '.md', '.pptx'}
    
    def __init__(
        self,
        chunk_size: int = settings.chunk_size,
        chunk_overlap: int = settings.chunk_overlap,
        session_id: str = None
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = ["\n\n", "\n", ". ", " ", ""]
        self.uploaded_files = []  # Track uploaded files for cleanup
        self.session_id = session_id  # For user-scoped storage
        
    def is_supported(self, filename: str) -> bool:
        """Check if file type is supported"""
        return Path(filename).suffix.lower() in self.SUPPORTED_EXTENSIONS
    
    def _load_pdf(self, file_path: Path) -> List[Document]:
        """Load PDF document"""
        documents = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        documents.append(Document(
                            page_content=text,
                            metadata={
                                'source': file_path.name,
                                'page': page_num + 1,
                                'file_type': '.pdf'
                            }
                        ))
        except Exception as e:
            log.error(f"Error reading PDF {file_path.name}: {str(e)}")
            raise
        return documents
    
    def _load_docx(self, file_path: Path) -> List[Document]:
        """Load DOCX document"""
        try:
            doc = docx.Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
            if not text.strip():
                raise ValueError("Document contains no text")
            return [Document(
                page_content=text,
                metadata={
                    'source': file_path.name,
                    'file_type': '.docx'
                }
            )]
        except Exception as e:
            log.error(f"Error reading DOCX {file_path.name}: {str(e)}")
            raise
    
    def _load_txt(self, file_path: Path) -> List[Document]:
        """Load text file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
            if not text.strip():
                raise ValueError("File is empty")
            return [Document(
                page_content=text,
                metadata={
                    'source': file_path.name,
                    'file_type': '.txt'
                }
            )]
        except Exception as e:
            log.error(f"Error reading TXT {file_path.name}: {str(e)}")
            raise
    
    def _load_markdown(self, file_path: Path) -> List[Document]:
        """Load Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
            if not text.strip():
                raise ValueError("File is empty")
            return [Document(
                page_content=text,
                metadata={
                    'source': file_path.name,
                    'file_type': '.md'
                }
            )]
        except Exception as e:
            log.error(f"Error reading MD {file_path.name}: {str(e)}")
            raise
    
    def _load_pptx(self, file_path: Path) -> List[Document]:
        """Load PowerPoint file"""
        try:
            prs = Presentation(file_path)
            documents = []
            
            for slide_num, slide in enumerate(prs.slides):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                
                text = "\n".join(text_parts)
                if text.strip():
                    documents.append(Document(
                        page_content=text,
                        metadata={
                            'source': file_path.name,
                            'slide': slide_num + 1,
                            'file_type': '.pptx'
                        }
                    ))
            
            if not documents:
                raise ValueError("No text content found in presentation")
            return documents
        except Exception as e:
            log.error(f"Error reading PPTX {file_path.name}: {str(e)}")
            raise
    
    def load_document(self, file_path: str) -> List[Document]:
        """Load a single document"""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = path.suffix.lower()
        if extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {extension}")
        
        log.info(f"Loading document: {path.name}")
        
        # Route to appropriate loader
        if extension == '.pdf':
            documents = self._load_pdf(path)
        elif extension == '.docx':
            documents = self._load_docx(path)
        elif extension == '.txt':
            documents = self._load_txt(path)
        elif extension == '.md':
            documents = self._load_markdown(path)
        elif extension == '.pptx':
            documents = self._load_pptx(path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")
        
        log.info(f"Loaded {len(documents)} page(s) from {path.name}")
        return documents
    
    def _split_text(self, text: str) -> List[str]:
        """Split text into chunks with overlap using recursive splitting"""
        if not text or not text.strip():
            return []
        
        def split_recursive(text: str, separators: List[str]) -> List[str]:
            """Recursively split text by separators"""
            if not separators:
                # No more separators, split by character count
                return [text[i:i + self.chunk_size] 
                       for i in range(0, len(text), max(1, self.chunk_size - self.chunk_overlap))]
            
            separator = separators[0]
            remaining_separators = separators[1:]
            
            if separator == "":
                # Character-level split
                return [text[i:i + self.chunk_size] 
                       for i in range(0, len(text), max(1, self.chunk_size - self.chunk_overlap))]
            
            splits = text.split(separator)
            result_chunks = []
            current_chunk = ""
            
            for split in splits:
                # If single split is too large, recursively split it
                if len(split) > self.chunk_size:
                    if current_chunk:
                        result_chunks.append(current_chunk.strip())
                        current_chunk = ""
                    # Recursively split large pieces
                    sub_chunks = split_recursive(split, remaining_separators)
                    result_chunks.extend(sub_chunks)
                elif len(current_chunk) + len(split) + len(separator) <= self.chunk_size:
                    current_chunk += split + separator
                else:
                    if current_chunk.strip():
                        result_chunks.append(current_chunk.strip())
                    current_chunk = split + separator
            
            if current_chunk.strip():
                result_chunks.append(current_chunk.strip())
            
            return result_chunks
        
        chunks = split_recursive(text, self.separators)
        
        # Apply overlap between chunks
        if self.chunk_overlap > 0 and len(chunks) > 1:
            final_chunks = []
            for i, chunk in enumerate(chunks):
                if i > 0 and final_chunks:
                    # Add overlap from previous chunk
                    overlap = final_chunks[-1][-self.chunk_overlap:] if len(final_chunks[-1]) > self.chunk_overlap else final_chunks[-1]
                    if overlap.strip():
                        chunk = overlap + " " + chunk
                final_chunks.append(chunk)
            return final_chunks
        
        return [c for c in chunks if c.strip()]
    
    def chunk_documents(self, documents: List[Document]) -> List[Document]:
        """Split documents into chunks"""
        log.info(f"Chunking {len(documents)} document(s)")
        
        all_chunks = []
        chunk_id = 0
        
        for doc in documents:
            text_chunks = self._split_text(doc.page_content)
            
            for chunk_text in text_chunks:
                if chunk_text.strip():
                    chunk_metadata = doc.metadata.copy()
                    chunk_metadata['chunk_id'] = chunk_id
                    
                    all_chunks.append(Document(
                        page_content=chunk_text,
                        metadata=chunk_metadata
                    ))
                    chunk_id += 1
        
        log.info(f"Created {len(all_chunks)} chunks")
        return all_chunks
    
    def process_file(self, file_path: str) -> List[Document]:
        """Load and chunk a single file"""
        documents = self.load_document(file_path)
        return self.chunk_documents(documents)
    
    def process_multiple_files(self, file_paths: List[str]) -> List[Document]:
        """Process multiple files and combine chunks"""
        all_chunks = []
        
        for file_path in file_paths:
            try:
                chunks = self.process_file(file_path)
                all_chunks.extend(chunks)
            except Exception as e:
                log.error(f"Failed to process {file_path}: {str(e)}")
                continue
        
        log.info(f"Total chunks from all files: {len(all_chunks)}")
        return all_chunks
    
    def save_uploaded_file(self, uploaded_file) -> str:
        """Save uploaded file to user-scoped uploads directory"""
        from config import UPLOADS_DIR, get_user_uploads_dir
        
        # Use session-scoped directory if session_id is set
        if self.session_id:
            upload_dir = get_user_uploads_dir(self.session_id)
        else:
            upload_dir = UPLOADS_DIR
        
        file_path = upload_dir / uploaded_file.name
        
        with open(file_path, 'wb') as f:
            f.write(uploaded_file.getbuffer())
        
        self.uploaded_files.append(str(file_path))
        log.info(f"Saved uploaded file: {file_path}")
        return str(file_path)
    
    def cleanup_uploaded_files(self):
        """Delete all uploaded files to free up disk space"""
        deleted_count = 0
        for file_path in self.uploaded_files:
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    deleted_count += 1
                    log.debug(f"Deleted: {file_path}")
            except Exception as e:
                log.warning(f"Failed to delete {file_path}: {str(e)}")
        
        self.uploaded_files.clear()
        if deleted_count > 0:
            log.info(f"Cleaned up {deleted_count} uploaded file(s)")
        
        return deleted_count
